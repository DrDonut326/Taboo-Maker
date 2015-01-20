import time
from pptx import Presentation
from pptx.util import Inches
import urllib
import urllib2
import os
import re
import sys

def error_line():
	print "!" * 25

def fancy_line():
	print " ^-^ " * 10

def housekeeping():
	print "Does the Pictures folder exist?"
	print os.path.isdir("Pictures")
	if not os.path.isdir("Pictures"):
		print "Don't worry <3 Making the folder now."
		os.makedirs("Pictures")
	print "Does input_file.txt exist?"
	print os.path.isfile("input_file.txt")
	if not os.path.isfile("input_file.txt"):
		print "Don't worry <3 Making the file now."
		open("input_file.txt", 'a')

def create_title_slide():
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]
	title.text = "Taboo Game"
	subtitle.text = "Don't say the word!"

def save_presentation():
	fancy_line()
	print "Type the name of your finished Presentation."    #save the presentation
	final_save_name = raw_input(">>> ")
	while not re.match("^[A-Za-z0-9_-]*$", final_save_name):   #check if the file name is a valid one
		print "Please use only characters or numbers with no spaces for the file name."
		final_save_name = raw_input(">>> ")
	prs.save(final_save_name+".pptx")
	print "Presentation Saved!  All done!"
	fancy_line()

def read_txt_file():
	print "Reading the input file."			#Read the txt file to get search terms.
	filename = "input_file.txt"
	txt = open(filename)
	print "Here are the words in the list: "+"\n"
	print txt.read()
	print "\n"
	raw_input("Press enter to continue or CNTRL-C to quit")  #Check to see if the user wants to continue or break
	txt.seek(0)  #Reset the cursor
	word_list = []
	while True:								#MAIN WHILE LOOP
											#Break the loop when there are no lines left to read
		line = txt.readline()
		if line == "": break
		stripped_line = line.rstrip('\n')
		word_list.append(stripped_line)
	txt.close()
	return word_list

def get_a_new_word(lst):
	"""Returns the first element of a list and pops that word off the list"""
	
	word = lst.pop(0)  #pops off the first word from the list
	word = re.sub("[^A-Za-z0-9]", "", word)  #strips out everything except letters and numbers using regEx
	return word,lst   #returns the cleaned word and the shortened list

def get_a_potential_url(search_term):
	print "Beginning URL finding."
	print search_term
	jpgto_search_url = "http://"+search_term + ".jpg.to/"
	#get the html for the website
	response = urllib2.urlopen(jpgto_search_url)
	print "Html file read."
	html = response.read()
	print "Html saved.  Starting parse."
	#set the index search term for finding the picture_url
	src_index = html.index('src=')
	#set beginning and end points of the parsing thingy
	startp = html.find('"',src_index)
	endp = html.rfind('"')
	#get the actual picture url
	picture_url = html[startp+1:endp]
	print "URL parsed."
	return picture_url	
	
def url_tester(url):
	"""Returns True or False if the url works or not."""
	try:
		resp = urllib2.urlopen(url) 	#try to open the url 
		print "URL good!"
		return True
	except urllib2.URLError:			#see if the url itself is broken
		print "URL BROKEN!"
		return False
	
	#except BadStatusLine(line):
	#	print "Something crazy with the URL."
	#	return False
	except:								#check for unexpected errors
		print "Unexpected error in url_tester():", sys.exc_info()[0]
		return False
		#raise
	
def get_working_url(word):
	"""Takes an initial word and returns a matching+working url impage for that word."""
	possible_url = get_a_potential_url(word)
	while True:
		test_result = url_tester(possible_url)
		if test_result == True:
			break
		else:
			if word.endswith("111"):
				print "Can't find a picture, changing to error picture."
				word = "erroricon"
			else:
				word = word + "1"
			possible_url = get_a_potential_url(word)
	return possible_url

def save_the_picture(url,search_term):
	print "Saving the picture, please wait..."
	save_name = search_term+".jpg"
	try:
		urllib.urlretrieve(final_url, "Pictures\\"+save_name)
		print "%s picture saved!" % search_term
	except urllib.IOError:
		error_line()
		print "Some kind of filter has prevented the program from downloading the picture."
		print "(probably the education office filter)"
		print "Downloading an error picture instead."
		error_line()
		urllib.urlretrieve("http://rationalwiki.org/w/images/f/f5/Error_icon.svg", "Pictures\\"+save_name)
	except:
		print "OH GOD EVERYTHING IS BROKEN :( "
		raise
		

def add_picture_to_slideshow(pic_name):
	#formats a title and content slide and adds it
	slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
	slide = prs.slides.add_slide(slide_layout)
	title = slide.shapes.title
	title.text = pic_name

	#adds a picture to the slide
	file_pic_name = "Pictures\\%s.jpg" % pic_name
	img_path = file_pic_name

	#sets size of image and adds it
	left = Inches(.4)
	top = Inches(1.5)
	width = Inches(9.25)
	height = Inches(5.75)
	try:
		pic = slide.shapes.add_picture(img_path, left, top, width, height)
	except IOError:
		error_line()
		print "Error adding the file to the ppt."
		print "Probably due to some kind of content filter."
		print "Adding skipping this picture."
		error_line()
	except:
		print "Unexpected error when trying to add the picture to a slide:", sys.exc_info()[0]
		raise

def star_line():
	print "*" * 50

def test_break():
	raw_input("TEST! BREAK NOW!")
	
star_line()
print "Taboo Game Maker."
print "Created by Patrick Naven, 2014."
star_line()


	
SLD_LAYOUT_TITLE_AND_CONTENT = 1
prs = Presentation()	
	
housekeeping() #checks / creates the necessary file and folder  #working!
create_title_slide() #make the title slide 
list_of_words = read_txt_file() #return an array containing list of words

while len(list_of_words)>0:
	
	star_line()
	new_word,list_of_words = get_a_new_word(list_of_words)
	print "Getting a picture for "+new_word
	final_url = get_working_url(new_word)
	save_the_picture(final_url,new_word)
	add_picture_to_slideshow(new_word)
	star_line()
	



save_presentation()


